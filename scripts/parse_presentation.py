"""Převede .pptx na plain text se slide tituly pro analýzu transcriptů."""

import sys
from pathlib import Path


def parse_pptx(path: Path) -> str:
    try:
        from pptx import Presentation
    except ImportError:
        print("Chybí závislost: pip install python-pptx", file=sys.stderr)
        sys.exit(1)

    prs = Presentation(path)
    lines = []
    for i, slide in enumerate(prs.slides, start=1):
        title = ""
        body_parts = []
        for shape in slide.shapes:
            if not shape.has_text_frame:
                continue
            text = shape.text_frame.text.strip()
            if not text:
                continue
            if shape.shape_id == slide.shapes.title.shape_id if slide.shapes.title else False:
                title = text
            else:
                body_parts.append(text)

        header = f"## Slide {i}" + (f" — {title}" if title else "")
        lines.append(header)
        if body_parts:
            lines.extend(body_parts)
        lines.append("")

    return "\n".join(lines)


if __name__ == "__main__":
    if len(sys.argv) != 2:
        print(f"Použití: python {sys.argv[0]} presentations/lekce-XX.pptx", file=sys.stderr)
        sys.exit(1)

    input_path = Path(sys.argv[1])
    if not input_path.exists():
        print(f"Soubor nenalezen: {input_path}", file=sys.stderr)
        sys.exit(1)

    print(parse_pptx(input_path))
